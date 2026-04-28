"""
Generate the AI-Powered Resume Shortlisting System PPT
Run: python presentation/generate_ppt.py
Output: presentation/AI_Resume_Shortlisting_System.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── colour palette ────────────────────────────────────────────────────────────
PURPLE      = RGBColor(0x6C, 0x63, 0xFF)   # #6c63ff
DARK_BG     = RGBColor(0x12, 0x10, 0x2A)   # #12102A  (dark navy)
LIGHT_BG    = RGBColor(0x1E, 0x1B, 0x4B)   # #1e1b4b  (slightly lighter)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW      = RGBColor(0xFF, 0xD7, 0x00)
CYAN        = RGBColor(0x00, 0xE5, 0xFF)
GREEN       = RGBColor(0x00, 0xFF, 0x9D)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED_COL     = RGBColor(0xFF, 0x4D, 0x4D)
LIGHT_PURPLE= RGBColor(0xC4, 0xB5, 0xFD)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────────────

def add_gradient_bg(slide, c1=DARK_BG, c2=LIGHT_BG):
    """Fill slide background with a solid dark colour (gradient via shape overlay)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = c1


def add_rect(slide, left, top, width, height, fill_color=None, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.FREEFORM if False else 1,
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_title_area(slide, title, subtitle=None,
                   title_size=40, sub_size=22):
    """Centered title + optional subtitle."""
    # decorative accent bar
    bar = add_rect(slide,
                   Inches(0.5), Inches(0.3),
                   Inches(0.08), Inches(0.8),
                   fill_color=PURPLE)

    add_textbox(slide, title,
                Inches(0.7), Inches(0.2),
                Inches(11.8), Inches(1.0),
                font_size=title_size, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.7), Inches(1.05),
                    Inches(11.8), Inches(0.6),
                    font_size=sub_size, bold=False,
                    color=LIGHT_PURPLE, align=PP_ALIGN.LEFT)

    # horizontal divider
    line = slide.shapes.add_shape(1,
        Inches(0.5), Inches(1.75),
        Inches(12.3), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()


def bullet_slide(prs, title, bullets, sub=None, bullet_color=WHITE,
                 bullet_size=20, icon_char="▶"):
    """Generic bullet-point slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_title_area(slide, title, sub)

    y = Inches(2.1)
    for item in bullets:
        if isinstance(item, tuple):
            icon, text = item
        else:
            icon, text = icon_char, item

        # icon
        add_textbox(slide, icon,
                    Inches(0.5), y, Inches(0.5), Inches(0.55),
                    font_size=bullet_size, bold=True,
                    color=PURPLE)
        # text
        add_textbox(slide, text,
                    Inches(1.05), y, Inches(11.5), Inches(0.55),
                    font_size=bullet_size, bold=False,
                    color=bullet_color)
        y += Inches(0.62)
    return slide


def table_slide(prs, title, headers, rows, col_widths=None, sub=None):
    """Slide with a formatted table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_title_area(slide, title, sub)

    n_cols = len(headers)
    n_rows = len(rows)
    if col_widths is None:
        col_widths = [Inches(12.3 / n_cols)] * n_cols

    tbl_left = Inches(0.5)
    tbl_top  = Inches(2.1)
    tbl_h    = Inches(0.5 * (n_rows + 1))

    table = slide.shapes.add_table(
        n_rows + 1, n_cols, tbl_left, tbl_top,
        sum(col_widths), tbl_h
    ).table

    # header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER

    # data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            bg = RGBColor(0x2D, 0x2A, 0x6E) if r % 2 == 0 else RGBColor(0x1E, 0x1B, 0x4B)
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.color.rgb = WHITE
            p.runs[0].font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT

    # set column widths
    tbl_xml = table._tbl
    tblGrid = tbl_xml.find(nsmap.qn('a:tblGrid'))
    for i, gridCol in enumerate(tblGrid.findall(nsmap.qn('a:gridCol'))):
        gridCol.set('w', str(col_widths[i]))

    return slide


# ── BUILD PRESENTATION ────────────────────────────────────────────────────────

def build_ppt(out_path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide1, DARK_BG, DARK_BG)

    # big purple accent block left
    add_rect(slide1, Inches(0), Inches(0), Inches(0.18), SLIDE_H, PURPLE)

    # decorative circle top-right
    circ = slide1.shapes.add_shape(9,    # oval
        Inches(10.8), Inches(-0.5), Inches(3), Inches(3))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x2D, 0x1B, 0x69)
    circ.line.fill.background()

    # title
    add_textbox(slide1, "AI-Powered Resume\nShortlisting System",
                Inches(0.5), Inches(1.5), Inches(9), Inches(2.4),
                font_size=44, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    # subtitle
    add_textbox(slide1,
                "Automating Candidate Screening Using NLP & Machine Learning",
                Inches(0.5), Inches(3.8), Inches(9), Inches(0.7),
                font_size=22, bold=False, color=LIGHT_PURPLE,
                align=PP_ALIGN.LEFT)

    # divider
    div = slide1.shapes.add_shape(1, Inches(0.5), Inches(4.55),
                                   Inches(5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = PURPLE
    div.line.fill.background()

    # meta info
    add_textbox(slide1,
                "Major Project Presentation\nDepartment of Computer Science & Engineering",
                Inches(0.5), Inches(4.8), Inches(8), Inches(0.9),
                font_size=16, bold=False, color=LIGHT_PURPLE)

    # tag line bottom
    add_textbox(slide1, "BERT  ·  TF-IDF  ·  React  ·  Node.js  ·  Flask  ·  Docker",
                Inches(0.5), Inches(6.6), Inches(10), Inches(0.55),
                font_size=14, color=RGBColor(0x94, 0x8F, 0xD4))

    # ── Slide 2: Problem Statement ────────────────────────────────────────────
    bullet_slide(prs,
        "Problem Statement",
        [
            ("⏱", "Recruiters spend only 6–8 seconds per resume — screening is rushed"),
            ("📄", "High application volumes cause fatigue and inconsistency"),
            ("⚠", "Qualified candidates are overlooked due to subjective screening"),
            ("🎯", "Bias creeps in through manual review of hundreds of resumes"),
            ("✅", "Need: An intelligent, automated, unbiased screening system"),
        ],
        bullet_size=20)

    # ── Slide 3: Objectives ───────────────────────────────────────────────────
    bullet_slide(prs,
        "Project Objectives",
        [
            ("1️", "Automate resume parsing from PDF and DOCX formats"),
            ("2️", "Match candidate skills against job requirements using AI"),
            ("3️", "Rank candidates objectively based on multiple weighted factors"),
            ("4️", "Provide recruiters a real-time visual dashboard"),
            ("5️", "Reduce hiring bias and improve decision speed"),
            ("6️", "Build a production-ready microservices architecture"),
        ],
        bullet_size=20)

    # ── Slide 4: System Architecture ─────────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide4)
    add_title_area(slide4, "System Architecture",
                   "Three-Tier Microservices Design")

    boxes = [
        (Inches(0.6),  Inches(2.4), "Frontend\n(React)\n:3000",
         RGBColor(0x4C, 0x1D, 0x95)),
        (Inches(5.05), Inches(2.4), "Backend\n(Express)\n:5000",
         RGBColor(0x1E, 0x40, 0xAF)),
        (Inches(9.5),  Inches(2.4), "ML Pipeline\n(Flask)\n:5001",
         RGBColor(0x06, 0x5F, 0x46)),
    ]
    box_w = Inches(3.2)
    box_h = Inches(1.8)

    for bx, by, txt, col in boxes:
        b = slide4.shapes.add_shape(1, bx, by, box_w, box_h)
        b.fill.solid(); b.fill.fore_color.rgb = col
        b.line.color.rgb = PURPLE; b.line.width = Pt(2)
        tf = b.text_frame
        tf.text = txt
        tf.paragraphs[0].runs[0].font.size = Pt(18)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # arrows between boxes
    for ax in [Inches(3.8), Inches(8.25)]:
        arr = slide4.shapes.add_shape(1, ax, Inches(3.1), Inches(1.25), Pt(4))
        arr.fill.solid(); arr.fill.fore_color.rgb = PURPLE
        arr.line.fill.background()
        add_textbox(slide4, "◄──►", ax, Inches(3.0), Inches(1.25), Inches(0.5),
                    font_size=20, color=PURPLE, align=PP_ALIGN.CENTER)

    # data flow steps
    steps = [
        "1. Recruiter posts job with required skills",
        "2. Candidates upload PDF/DOCX resumes",
        "3. ML Pipeline parses resume → raw text",
        "4. BERT + TF-IDF calculate semantic similarity",
        "5. Weighted score → Ranked dashboard",
    ]
    y = Inches(4.5)
    for s in steps:
        add_textbox(slide4, s, Inches(0.6), y, Inches(12), Inches(0.45),
                    font_size=16, color=LIGHT_PURPLE)
        y += Inches(0.46)

    # ── Slide 5: Tech Stack ───────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide5)
    add_title_area(slide5, "Tech Stack")

    categories = [
        ("Frontend",     PURPLE,                   ["React 18", "React Router", "Axios", "CSS3"]),
        ("Backend",      RGBColor(0x1E,0x40,0xAF), ["Node.js", "Express.js", "Multer", "JWT"]),
        ("ML Pipeline",  RGBColor(0x06,0x5F,0x46), ["Python 3", "Flask", "scikit-learn",
                                                      "Sentence-BERT", "PyPDF2", "python-docx"]),
        ("DevOps",       RGBColor(0x92,0x40,0x0E), ["Docker", "Docker Compose", "PowerShell"]),
    ]

    col_w = Inches(3.0)
    col_h = Inches(4.5)
    gap   = Inches(0.2)
    start_x = Inches(0.35)

    for i, (cat, col, items) in enumerate(categories):
        cx = start_x + i * (col_w + gap)
        # header box
        hdr = slide5.shapes.add_shape(1, cx, Inches(2.0), col_w, Inches(0.65))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = col
        hdr.line.fill.background()
        tf = hdr.text_frame
        tf.text = cat
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = WHITE
        tf.paragraphs[0].runs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # body box
        body = slide5.shapes.add_shape(1, cx, Inches(2.65), col_w, col_h)
        body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0x2D,0x2A,0x6E)
        body.line.color.rgb = col; body.line.width = Pt(1.5)

        ty = Inches(2.75)
        for item in items:
            add_textbox(slide5, f"• {item}", cx + Inches(0.1), ty,
                        col_w - Inches(0.15), Inches(0.48),
                        font_size=15, color=WHITE)
            ty += Inches(0.52)

    # ── Slide 6: Data Flow ────────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide6)
    add_title_area(slide6, "Data Flow", "End-to-End Processing Pipeline")

    steps6 = [
        ("1", "Recruiter posts job with title, description & required skills",     PURPLE),
        ("2", "Candidate uploads PDF/DOCX resume via the React UI",                RGBColor(0x1E,0x40,0xAF)),
        ("3", "ML Pipeline parses resume → extracts raw text (PyPDF2 / python-docx)", RGBColor(0x06,0x5F,0x46)),
        ("4", "BERT + TF-IDF compute semantic similarity between resume & JD",      RGBColor(0x78,0x35,0x09)),
        ("5", "Skill-matching engine detects required skills with synonym support",  RGBColor(0x70,0x1A,0x75)),
        ("6", "Combined weighted score computed (Skills 65% + BERT 12% + …)",       RGBColor(0x1E,0x3A,0x8A)),
        ("7", "Ranked candidate list displayed in real-time recruiter dashboard",    RGBColor(0x06,0x4E,0x3B)),
    ]

    y = Inches(2.0)
    for num, text, col in steps6:
        # number badge
        badge = slide6.shapes.add_shape(9, Inches(0.5), y, Inches(0.55), Inches(0.55))
        badge.fill.solid(); badge.fill.fore_color.rgb = col
        badge.line.fill.background()
        tf = badge.text_frame
        tf.text = num
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = WHITE
        tf.paragraphs[0].runs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide6, text, Inches(1.2), y, Inches(11.5), Inches(0.55),
                    font_size=18, color=WHITE)

        # arrow connector
        if num != "7":
            arr = slide6.shapes.add_shape(1,
                Inches(0.68), y + Inches(0.55), Pt(4), Inches(0.22))
            arr.fill.solid(); arr.fill.fore_color.rgb = col
            arr.line.fill.background()

        y += Inches(0.76)

    # ── Slide 7: ML Algorithm ─────────────────────────────────────────────────
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide7)
    add_title_area(slide7, "ML Scoring Algorithm",
                   "Weighted Multi-Factor Formula")

    components = [
        ("Skills Match",         65, PURPLE,                  "Synonym detection (React=React.js, DSA=Data Structures)"),
        ("BERT Similarity",      12, RGBColor(0x1E,0x40,0xAF),"Sentence-BERT all-MiniLM-L6-v2 — semantic context"),
        ("Experience/Education", 10, RGBColor(0x06,0x5F,0x46),"Years of exp (0-2, 2-5, 5+) + degree level"),
        ("Project Relevance",     8, RGBColor(0x78,0x35,0x09),"Keyword analysis of projects section"),
        ("Role Keywords",         3, RGBColor(0x70,0x1A,0x75),"Seniority levels: Fresher, Junior, Senior, Lead"),
        ("Certifications",        2, RGBColor(0x1E,0x3A,0x8A),"AWS, Azure, MongoDB, GCP, etc."),
    ]

    bar_left = Inches(0.5)
    bar_top  = Inches(2.1)
    max_bar  = Inches(8.5)
    row_h    = Inches(0.72)

    for i, (label, pct, col, detail) in enumerate(components):
        y = bar_top + i * row_h
        # label
        add_textbox(slide7, f"{label}", bar_left, y, Inches(2.5), Inches(0.5),
                    font_size=15, bold=True, color=WHITE)
        # bar
        bw = max_bar * pct / 100
        b = slide7.shapes.add_shape(1, bar_left + Inches(2.5), y + Inches(0.08),
                                     bw, Inches(0.38))
        b.fill.solid(); b.fill.fore_color.rgb = col
        b.line.fill.background()
        # pct label
        add_textbox(slide7, f"{pct}%",
                    bar_left + Inches(2.5) + bw + Inches(0.08), y,
                    Inches(0.5), Inches(0.5),
                    font_size=15, bold=True, color=col)
        # detail small
        add_textbox(slide7, detail,
                    bar_left + Inches(2.5), y + Inches(0.40),
                    Inches(9.5), Inches(0.28),
                    font_size=11, color=LIGHT_PURPLE)

    # formula box at bottom
    formula_box = slide7.shapes.add_shape(1, Inches(0.5), Inches(6.8),
                                           Inches(12.3), Inches(0.58))
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = RGBColor(0x1E,0x1B,0x4B)
    formula_box.line.color.rgb = PURPLE; formula_box.line.width = Pt(1.5)
    add_textbox(slide7,
        "Score = (Skills×0.65) + (BERT×0.12) + (Exp×0.10) + (Projects×0.08) + (Keywords×0.03) + (Certs×0.02)",
        Inches(0.6), Inches(6.82), Inches(12.0), Inches(0.52),
        font_size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # ── Slide 8: ML Techniques ────────────────────────────────────────────────
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide8)
    add_title_area(slide8, "Key ML Techniques Explained")

    techniques = [
        ("TF-IDF Vectorization",
         "Converts resume/JD text to numerical vectors.\nWeights rare but important terms higher (e.g., specific tech skills).",
         PURPLE),
        ("Cosine Similarity",
         "Measures the angle between two vectors.\n0 = no match, 1 = perfect match — direction matters, not magnitude.",
         RGBColor(0x1E,0x40,0xAF)),
        ("Sentence-BERT",
         "Pre-trained transformer model (all-MiniLM-L6-v2).\nUnderstands meaning — 'software engineer' ≈ 'developer'.",
         RGBColor(0x06,0x5F,0x46)),
        ("N-gram Skill Matching",
         "Handles multi-word skills: 'machine learning', 'data structures'.\nSynonym mapping covers abbreviations (DSA = Data Structures).",
         RGBColor(0x78,0x35,0x09)),
    ]

    cols = 2
    card_w = Inches(6.0)
    card_h = Inches(2.0)
    xs = [Inches(0.4), Inches(6.85)]
    ys = [Inches(2.1), Inches(4.4)]

    for idx, (title_t, desc, col) in enumerate(techniques):
        cx = xs[idx % cols]
        cy = ys[idx // cols]
        card = slide8.shapes.add_shape(1, cx, cy, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x1E,0x1B,0x4B)
        card.line.color.rgb = col; card.line.width = Pt(2)
        # accent top strip
        strip = slide8.shapes.add_shape(1, cx, cy, card_w, Inches(0.08))
        strip.fill.solid(); strip.fill.fore_color.rgb = col
        strip.line.fill.background()
        add_textbox(slide8, title_t, cx + Inches(0.15), cy + Inches(0.12),
                    card_w - Inches(0.3), Inches(0.45),
                    font_size=17, bold=True, color=col)
        add_textbox(slide8, desc, cx + Inches(0.15), cy + Inches(0.6),
                    card_w - Inches(0.3), Inches(1.3),
                    font_size=14, color=WHITE)

    # ── Slide 9: Resume Parsing ───────────────────────────────────────────────
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide9)
    add_title_area(slide9, "Resume Parsing Module",
                   "Automated extraction from PDF & DOCX")

    parsing_points = [
        ("📄", "Supports PDF (via PyPDF2) and DOCX (via python-docx)"),
        ("🔍", "Extracts: Name, Email, Skills, Projects, Experience, Education, Certifications"),
        ("🧩", "Regex-based section detection with intelligent fallback for unstructured resumes"),
        ("🔒", "All parsing happens server-side — no data sent to third parties"),
        ("🔄", "Parsed text sent to similarity engine for scoring"),
    ]

    y = Inches(2.2)
    for icon, text in parsing_points:
        add_textbox(slide9, icon, Inches(0.5), y, Inches(0.7), Inches(0.6),
                    font_size=24)
        add_textbox(slide9, text, Inches(1.3), y, Inches(11.3), Inches(0.6),
                    font_size=19, color=WHITE)
        y += Inches(0.75)

    # flow box
    flow_box = slide9.shapes.add_shape(1, Inches(0.5), Inches(6.5),
                                        Inches(12.3), Inches(0.75))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(0x2D,0x2A,0x6E)
    flow_box.line.color.rgb = PURPLE; flow_box.line.width = Pt(1.5)
    add_textbox(slide9,
        "Resume File  →  Text Extraction  →  Section Detection  →  Structured Data  →  ML Scoring",
        Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.6),
        font_size=15, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # ── Slide 10: API Design ──────────────────────────────────────────────────
    table_slide(prs,
        "RESTful API Design",
        ["Endpoint", "Method", "Description"],
        [
            ["/api/jobs",                   "POST / GET", "Create & list job postings"],
            ["/api/jobs/:id",               "GET / PUT / DELETE", "Manage individual job"],
            ["/api/resumes/upload",         "POST", "Upload resume (multipart/form-data)"],
            ["/api/resumes/match-jobs",     "POST", "Match resume against all jobs"],
            ["/api/rankings/generate/:jobId","POST", "Generate AI rankings for a job"],
            ["/api/rankings/job/:jobId",    "GET",  "Retrieve rankings for a job"],
            ["/parse/resume",               "POST", "ML service: parse PDF/DOCX"],
            ["/ranking/calculate",          "POST", "ML service: calculate scores"],
        ],
        col_widths=[Inches(3.8), Inches(2.5), Inches(6.0)],
        sub="8 endpoints · JSON responses · { success, data } format"
    )

    # ── Slide 11: Frontend Features ───────────────────────────────────────────
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide11)
    add_title_area(slide11, "Frontend Features",
                   "React 18 · Purple Gradient UI · Tabbed Navigation")

    features = [
        ("Post Job Tab",      PURPLE,                   [
            "Job title, company, description",
            "Add required skills (comma-separated)",
            "Form validation & success feedback",
        ]),
        ("Upload Resume Tab", RGBColor(0x1E,0x40,0xAF), [
            "Candidate name & email",
            "Select target job from dropdown",
            "PDF / DOCX file upload",
        ]),
        ("Dashboard Tab",     RGBColor(0x06,0x5F,0x46), [
            "Select job → Generate Rankings button",
            "Ranked list: match %, rank position",
            "Matched skills vs. missing skills",
        ]),
    ]

    feat_w = Inches(3.9)
    gap    = Inches(0.25)
    sx     = Inches(0.35)

    for i, (name, col, items) in enumerate(features):
        fx = sx + i * (feat_w + gap)
        fy = Inches(2.1)
        # header
        hdr = slide11.shapes.add_shape(1, fx, fy, feat_w, Inches(0.6))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = col
        hdr.line.fill.background()
        add_textbox(slide11, name, fx, fy + Inches(0.08), feat_w, Inches(0.45),
                    font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # body
        body = slide11.shapes.add_shape(1, fx, fy + Inches(0.6), feat_w, Inches(4.2))
        body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0x1E,0x1B,0x4B)
        body.line.color.rgb = col; body.line.width = Pt(1.5)
        ty = fy + Inches(0.75)
        for item in items:
            add_textbox(slide11, f"✓  {item}", fx + Inches(0.12), ty,
                        feat_w - Inches(0.2), Inches(0.55),
                        font_size=15, color=WHITE)
            ty += Inches(0.65)

    add_textbox(slide11,
        "Modern purple-gradient design  ·  Responsive layout  ·  Real-time success/error notifications",
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.38),
        font_size=13, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)

    # ── Slide 12: Score Interpretation ────────────────────────────────────────
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide12)
    add_title_area(slide12, "Score Interpretation",
                   "How to read the match percentage")

    ranges = [
        ("60 – 75%+", "✅ Strong Match",   "Definitely invite for interview",  GREEN,              Inches(0.6)),
        ("50 – 60%",  "🟡 Good Match",     "Seriously consider the candidate", YELLOW,             Inches(2.0)),
        ("40 – 50%",  "🟠 Decent Match",   "Review manually, could be a fit",  ORANGE,             Inches(3.4)),
        ("< 40%",     "🔴 Weak Fit",       "Likely not a match for this role", RED_COL,            Inches(4.8)),
    ]

    for score, badge, desc, col, by in ranges:
        # score band box
        sb = slide12.shapes.add_shape(1, Inches(0.5), by, Inches(2.2), Inches(1.1))
        sb.fill.solid(); sb.fill.fore_color.rgb = col
        sb.line.fill.background()
        add_textbox(slide12, score, Inches(0.5), by + Inches(0.25), Inches(2.2), Inches(0.55),
                    font_size=22, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        # badge
        add_textbox(slide12, badge, Inches(2.9), by + Inches(0.05),
                    Inches(4.0), Inches(0.55),
                    font_size=22, bold=True, color=col)
        # description
        add_textbox(slide12, desc, Inches(2.9), by + Inches(0.6),
                    Inches(9.5), Inches(0.45),
                    font_size=17, color=LIGHT_PURPLE)

    # ── Slide 13: Project Structure ───────────────────────────────────────────
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide13)
    add_title_area(slide13, "Project Structure",
                   "Clean separation of concerns across 5 modules")

    tree_lines = [
        ("MAJOR Project/",         YELLOW,       False),
        ("├── backend/",           LIGHT_PURPLE, False),
        ("│     Node.js + Express REST API, JWT auth, in-memory storage", WHITE, True),
        ("├── frontend/",          LIGHT_PURPLE, False),
        ("│     React 18 SPA — Post Job, Upload Resume, Rankings Dashboard", WHITE, True),
        ("├── ml-pipeline/",       LIGHT_PURPLE, False),
        ("│     Flask ML service: parser.py (PDF/DOCX) + similarity.py (BERT/TF-IDF)", WHITE, True),
        ("├── docker/",            LIGHT_PURPLE, False),
        ("│     docker-compose.yml + 3 Dockerfiles — one-command deployment", WHITE, True),
        ("└── START.ps1",          LIGHT_PURPLE, False),
        ("      One-click PowerShell startup script for local development", WHITE, True),
    ]

    y = Inches(2.1)
    for text, col, indent in tree_lines:
        x = Inches(1.0) if indent else Inches(0.5)
        sz = 15 if indent else 17
        bold = not indent
        add_textbox(slide13, text, x, y, Inches(12.0), Inches(0.42),
                    font_size=sz, bold=bold, color=col)
        y += Inches(0.44)

    # ── Slide 14: Deployment & DevOps ─────────────────────────────────────────
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide14)
    add_title_area(slide14, "Deployment & DevOps",
                   "Docker Compose · Microservices · One-Command Launch")

    # three service port cards
    ports = [
        ("Frontend",   ":3000", PURPLE),
        ("Backend",    ":5000", RGBColor(0x1E,0x40,0xAF)),
        ("ML Service", ":5001", RGBColor(0x06,0x5F,0x46)),
    ]
    cw = Inches(3.8); gap14 = Inches(0.2)
    sx14 = Inches(0.5)
    for i, (svc, port, col) in enumerate(ports):
        cx = sx14 + i * (cw + gap14)
        c = slide14.shapes.add_shape(1, cx, Inches(2.1), cw, Inches(1.2))
        c.fill.solid(); c.fill.fore_color.rgb = col
        c.line.fill.background()
        add_textbox(slide14, svc, cx, Inches(2.15), cw, Inches(0.55),
                    font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide14, port, cx, Inches(2.65), cw, Inches(0.55),
                    font_size=24, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    deploy_points = [
        ("🐳", "Docker Compose: docker-compose up --build  →  all 3 services start"),
        ("⚡", "One-click local dev: .\\START.ps1  (PowerShell)"),
        ("🔒", "Each service runs in an isolated container with its own network"),
        ("🚀", "Production path: MongoDB + AWS S3 + Vercel / Heroku / AWS ECS"),
        ("📦", "Each Dockerfile: backend, frontend, ml — independently scalable"),
    ]

    y14 = Inches(3.6)
    for icon, text in deploy_points:
        add_textbox(slide14, icon, Inches(0.5), y14, Inches(0.6), Inches(0.55), font_size=22)
        add_textbox(slide14, text, Inches(1.2), y14, Inches(11.5), Inches(0.55),
                    font_size=18, color=WHITE)
        y14 += Inches(0.68)

    # ── Slide 15: Challenges & Solutions ──────────────────────────────────────
    table_slide(prs,
        "Challenges & Solutions",
        ["Challenge", "Solution"],
        [
            ["Skill name variations (React vs React.js)", "Synonym mapping dictionary in skill matcher"],
            ["Unstructured resume formats",               "Regex section detection + plain-text fallback"],
            ["BERT model unavailable / slow load",        "Graceful fallback to TF-IDF cosine similarity"],
            ["Multi-service coordination",                "Docker Compose with health checks & restart policy"],
            ["Bias in manual hiring",                     "Objective, weighted, reproducible scoring formula"],
            ["High resume volume",                        "In-memory processing — no DB round-trips needed"],
        ],
        col_widths=[Inches(6.0), Inches(6.3)],
        sub="Engineering decisions that shaped the final architecture"
    )

    # ── Slide 16: Results & Demo ───────────────────────────────────────────────
    slide16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide16)
    add_title_area(slide16, "Results & Demo",
                   "Sample ranking output for: Senior Full Stack Developer")

    # workflow steps
    wf_steps = ["1. Post Job", "2. Upload 3-5 Resumes", "3. Generate Rankings"]
    swx = Inches(0.5)
    for step in wf_steps:
        wb = slide16.shapes.add_shape(1, swx, Inches(2.1), Inches(3.9), Inches(0.6))
        wb.fill.solid(); wb.fill.fore_color.rgb = RGBColor(0x2D,0x2A,0x6E)
        wb.line.color.rgb = PURPLE; wb.line.width = Pt(1.5)
        add_textbox(slide16, step, swx, Inches(2.15), Inches(3.9), Inches(0.5),
                    font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        swx += Inches(4.1)

    # sample ranked results
    candidates = [
        ("🥇 Rank 1", "Candidate A — React, Node.js, MongoDB, Docker",    "72%", GREEN),
        ("🥈 Rank 2", "Candidate B — React, Node.js, Python",              "58%", YELLOW),
        ("🥉 Rank 3", "Candidate C — Node.js, SQL (missing: React, Docker)","41%", ORANGE),
    ]

    y16 = Inches(3.0)
    for rank, name, score, col in candidates:
        rb = slide16.shapes.add_shape(1, Inches(0.5), y16, Inches(11.5), Inches(0.85))
        rb.fill.solid(); rb.fill.fore_color.rgb = RGBColor(0x1E,0x1B,0x4B)
        rb.line.color.rgb = col; rb.line.width = Pt(2)
        add_textbox(slide16, rank, Inches(0.6), y16 + Inches(0.17),
                    Inches(1.3), Inches(0.5), font_size=17, bold=True, color=col)
        add_textbox(slide16, name, Inches(1.9), y16 + Inches(0.17),
                    Inches(8.5), Inches(0.5), font_size=16, color=WHITE)
        add_textbox(slide16, score, Inches(11.0), y16 + Inches(0.1),
                    Inches(1.0), Inches(0.65), font_size=28, bold=True,
                    color=col, align=PP_ALIGN.CENTER)
        y16 += Inches(0.98)

    add_textbox(slide16,
        "Each candidate card shows: matched skills ✅ | missing skills ❌ | overall match percentage",
        Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
        font_size=14, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)

    # ── Slide 17: Learning Outcomes ───────────────────────────────────────────
    bullet_slide(prs,
        "Learning Outcomes",
        [
            ("✅", "Full-stack MERN development (MongoDB · Express · React · Node.js)"),
            ("✅", "Machine Learning & NLP integration in a real-world application"),
            ("✅", "Microservices architecture design and service decomposition"),
            ("✅", "RESTful API design best practices and clean JSON contracts"),
            ("✅", "File handling — parsing PDF & DOCX on the server side"),
            ("✅", "Docker containerisation and multi-service orchestration"),
            ("✅", "Real-time data processing without a persistent database"),
        ],
        bullet_size=19)

    # ── Slide 18: Future Enhancements ─────────────────────────────────────────
    slide18 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide18)
    add_title_area(slide18, "Future Enhancements",
                   "Planned features for production readiness")

    enhancements = [
        ("🔐", "Role-based authentication (JWT) — Recruiter vs. Candidate portals"),
        ("📊", "Advanced analytics dashboard — hiring trends, funnel metrics"),
        ("📧", "Automated email notifications to shortlisted candidates"),
        ("📱", "Mobile-responsive UI improvements & PWA support"),
        ("🤖", "LinkedIn profile integration & ATS (Applicant Tracking System)"),
        ("☁️",  "Full cloud deployment on AWS / Azure / GCP"),
        ("📈", "Bias detection module & diversity-aware scoring"),
        ("🧪", "Unit & integration test coverage (Jest + PyTest)"),
    ]

    cols18 = 2; col18_w = Inches(6.0); gap18 = Inches(0.33)
    xs18 = [Inches(0.35), Inches(6.68)]
    y18  = Inches(2.1)
    row18_h = Inches(0.65)

    for i, (icon, text) in enumerate(enhancements):
        cx = xs18[i % cols18]
        cy = y18 + (i // cols18) * row18_h
        add_textbox(slide18, icon, cx, cy, Inches(0.55), Inches(0.55), font_size=22)
        add_textbox(slide18, text, cx + Inches(0.6), cy, Inches(5.6), Inches(0.55),
                    font_size=16, color=WHITE)

    # ── Slide 19: Conclusion ───────────────────────────────────────────────────
    slide19 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide19)
    add_title_area(slide19, "Conclusion")

    conclusions = [
        "Built a complete, end-to-end AI-powered recruitment tool from scratch",
        "Reduces manual screening time from hours to under a minute",
        "Objective & consistent ranking eliminates subjective bias",
        "Demonstrates seamless integration of modern web tech with ML/NLP",
        "Production-ready microservices architecture with Docker support",
        "Extensible design — ready for authentication, analytics & cloud deployment",
    ]

    y19 = Inches(2.1)
    for c in conclusions:
        box = slide19.shapes.add_shape(1, Inches(0.5), y19, Inches(12.3), Inches(0.72))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x2D,0x2A,0x6E)
        box.line.color.rgb = PURPLE; box.line.width = Pt(1)
        add_textbox(slide19, f"  ★  {c}", Inches(0.55), y19 + Inches(0.1),
                    Inches(12.1), Inches(0.55), font_size=18, color=WHITE)
        y19 += Inches(0.84)

    # ── Slide 20: Thank You ────────────────────────────────────────────────────
    slide20 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide20, DARK_BG, DARK_BG)

    add_rect(slide20, Inches(0), Inches(0), Inches(0.18), SLIDE_H, PURPLE)

    add_textbox(slide20, "Thank You!",
                Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.4),
                font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide20, "Questions & Discussion",
                Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.7),
                font_size=26, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)

    # divider
    div20 = slide20.shapes.add_shape(1, Inches(3.0), Inches(3.75),
                                      Inches(7.3), Pt(3))
    div20.fill.solid(); div20.fill.fore_color.rgb = PURPLE
    div20.line.fill.background()

    refs = [
        "scikit-learn — Machine Learning in Python",
        "Sentence-BERT — Hugging Face Transformers",
        "React — Meta Open Source",
        "Express.js — Fast Node.js web framework",
        "Flask — Lightweight Python web framework",
        "PyPDF2 / python-docx — Document parsing",
    ]
    y20 = Inches(4.0)
    for ref in refs:
        add_textbox(slide20, f"• {ref}", Inches(2.5), y20, Inches(8.3), Inches(0.4),
                    font_size=15, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)
        y20 += Inches(0.42)

    add_textbox(slide20,
        "AI-Powered Resume Shortlisting System  |  Major Project  |  GitHub: AryanBartwal",
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.38),
        font_size=13, color=RGBColor(0x94,0x8F,0xD4), align=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(out_path)
    print(f"✅  Saved: {out_path}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    import os
    out = os.path.join(os.path.dirname(__file__),
                       "AI_Resume_Shortlisting_System.pptx")
    build_ppt(out)
